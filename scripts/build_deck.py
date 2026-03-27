"""
SiBoNi CXO Cockpit — PowerPoint Deck Generator
Run: python scripts/build_deck.py
Output: siboni-deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette (light / executive) ──────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W   = RGBColor(0xF8, 0xFA, 0xFC)   # #F8FAFC
SLATE50 = RGBColor(0xF1, 0xF5, 0xF9)   # very light blue-grey bg
SLATE1  = RGBColor(0x0F, 0x17, 0x2A)   # near-black text
SLATE2  = RGBColor(0x1E, 0x29, 0x3B)
SLATE3  = RGBColor(0x33, 0x4A, 0x5E)
BLUE    = RGBColor(0x1D, 0x4E, 0xD8)   # #1D4ED8
BLUE_LT = RGBColor(0xDB, 0xEA, 0xFE)   # #DBEAFE
INDIGO  = RGBColor(0x43, 0x38, 0xCA)
VIOLET  = RGBColor(0x6D, 0x28, 0xD9)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
AMBER_LT= RGBColor(0xFE, 0xF3, 0xC7)   # #FEF3C7
GREEN   = RGBColor(0x05, 0x96, 0x69)
GREEN_LT= RGBColor(0xD1, 0xFA, 0xE5)
RED     = RGBColor(0xDC, 0x26, 0x26)
RED_LT  = RGBColor(0xFE, 0xE2, 0xE2)
CYAN    = RGBColor(0x02, 0x79, 0x95)
VIOLET_LT=RGBColor(0xED, 0xE9, 0xFE)
GREY_BD = RGBColor(0xE2, 0xE8, 0xF0)   # border grey
GREY_TX = RGBColor(0x64, 0x74, 0x8B)   # muted text

def blank_slide():
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color=OFF_W):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill=WHITE, border=GREY_BD, border_pt=0.75, radius=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_pt:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=11, bold=False, color=SLATE1,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def pill_label(slide, text, x, y, fill, text_color, size=8):
    w = len(text) * 0.075 + 0.22
    shape = box(slide, x, y, w, 0.22, fill=fill, border=fill, border_pt=0)
    txb = slide.shapes.add_textbox(Inches(x+0.04), Inches(y+0.02), Inches(w-0.08), Inches(0.2))
    txb.word_wrap = False
    p = txb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return w

def header_band(slide, title, subtitle="", accent=BLUE):
    box(slide, 0, 0, 13.33, 1.05, fill=WHITE, border=GREY_BD, border_pt=0)
    # accent left bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.06), Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    txt(slide, title, 0.2, 0.08, 10, 0.42, size=22, bold=True, color=SLATE1)
    if subtitle:
        txt(slide, subtitle, 0.2, 0.55, 12, 0.38, size=11, color=GREY_TX)
    # page rule
    rule = slide.shapes.add_shape(1, Inches(0), Inches(1.05), Inches(13.33), Inches(0.02))
    rule.fill.solid(); rule.fill.fore_color.rgb = GREY_BD
    rule.line.fill.background()

def slide_number(slide, n, total=11):
    txt(slide, f"{n} / {total}", 12.5, 7.22, 0.8, 0.22,
        size=8, color=GREY_TX, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, WHITE)

# Top accent stripe
bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

# Logo / product name
txt(s, "SiBoNi", 0.5, 0.18, 3, 0.5, size=32, bold=True, color=BLUE)
txt(s, "CXO COCKPIT  ·  EXECUTIVE INTELLIGENCE PLATFORM", 0.5, 0.68, 10, 0.3,
    size=9, bold=True, color=GREY_TX)

# Headline
txt(s, "Decisions at the Speed of Intelligence", 0.5, 1.3, 12.3, 1.0,
    size=36, bold=True, color=SLATE1, align=PP_ALIGN.CENTER)

# Sub
txt(s,
    "SiBoNi connects every enterprise system — ERP, CRM, HR, IoT, market feeds — "
    "synthesises signals through always-on AI and delivers executive-grade intelligence "
    "to the CEO's morning view. In real time.",
    1.5, 2.45, 10.3, 0.9, size=13, color=GREY_TX, align=PP_ALIGN.CENTER)

# 4 pillar boxes
pillars = [
    ("9+ MCP\nConnectors",      BLUE_LT,    BLUE),
    ("AI-Synthesised\nInsights", VIOLET_LT, VIOLET),
    ("Real-Time\nDecision Support", GREEN_LT, GREEN),
    ("Analyst → CEO\nIntelligence Bridge", AMBER_LT, AMBER),
]
for i, (label, bg, fg) in enumerate(pillars):
    bx = box(s, 0.5 + i*3.1, 3.6, 2.9, 1.05, fill=bg, border=fg, border_pt=1)
    txt(s, label, 0.5 + i*3.1 + 0.15, 3.73, 2.6, 0.8,
        size=12, bold=True, color=fg, align=PP_ALIGN.CENTER)

# Tagline footer
txt(s, "Manufacturing · Industrial · Multi-plant · Multi-region  ·  Built on MCP Open Standard",
    0.5, 4.95, 12.3, 0.3, size=9, color=GREY_TX, align=PP_ALIGN.CENTER)

# Bottom rule
r2 = s.shapes.add_shape(1, Inches(0), Inches(7.28), Inches(13.33), Inches(0.02))
r2.fill.solid(); r2.fill.fore_color.rgb = BLUE
r2.line.fill.background()

slide_number(s, 1)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — E2E FLOW
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "End-to-End Flow", "From raw enterprise data to executive decision — in minutes", BLUE)

nodes = [
    ("Enterprise\nSources",      "SAP ERP · CRM\nWorkday · IoT\nReuters",       BLUE_LT,    BLUE),
    ("MCP\nConnectors",          "Unified API\nReal-time sync\nAuth & governance", VIOLET_LT, VIOLET),
    ("Normalise\n& Enrich",      "KPI compute\nAnomaly detect\nSignal scoring",  RGBColor(0xE0,0xF2,0xFE), CYAN),
    ("AI Synthesis\nEngine",     "System prompts\nAnalyst prompts\nConfidence score", AMBER_LT, AMBER),
    ("Decision\nIntelligence",   "Ranked decisions\nEvidence bundles\nImpact model",  GREEN_LT,   GREEN),
    ("CEO\nCockpit",             "Insights Hub\nDecision Hub\nBoard Brief",      RED_LT,     RED),
    ("Execution\nTracking",      "Action tracking\nKPI feedback\nBoard reporting",GREEN_LT,  GREEN),
]

node_w = 1.55
gap = 0.22
start_x = 0.3
for i, (title, detail, bg, fg) in enumerate(nodes):
    nx = start_x + i*(node_w + gap)
    bx = box(s, nx, 1.25, node_w, 1.7, fill=bg, border=fg, border_pt=1.25)
    txt(s, title, nx+0.07, 1.3, node_w-0.14, 0.45, size=10, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(s, detail, nx+0.07, 1.78, node_w-0.14, 1.1,  size=8, color=SLATE3, align=PP_ALIGN.CENTER)
    if i < len(nodes)-1:
        ax = nx + node_w + 0.03
        arr = s.shapes.add_shape(1, Inches(ax), Inches(1.93), Inches(gap-0.06), Inches(0.12))
        arr.fill.solid(); arr.fill.fore_color.rgb = GREY_TX
        arr.line.fill.background()

# Latency bar
box(s, 0.3, 3.2, 12.73, 0.72, fill=WHITE, border=GREY_BD)
txt(s, "LATENCY PROFILE", 0.55, 3.28, 2, 0.25, size=8, bold=True, color=GREY_TX)
latencies = [
    ("Data sync",      "1–15 min", BLUE),
    ("AI synthesis",   "< 2 sec",  AMBER),
    ("Analyst→CEO",    "Instant",  GREEN),
    ("Board pack",     "1-click",  VIOLET),
    ("Decision→Action","Same session", RED),
]
for i, (lbl, val, col) in enumerate(latencies):
    lx = 2.4 + i * 2.15
    txt(s, lbl, lx, 3.28, 2, 0.22, size=8, color=GREY_TX)
    txt(s, val, lx, 3.5,  2, 0.28, size=11, bold=True, color=col)

slide_number(s, 2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MCP CONNECTORS
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "MCP Connector Library", "Model Context Protocol — live context from every enterprise system", VIOLET)

connectors = [
    ("SAP S/4HANA ERP",     "Orders · Inventory · Production · GL",  "LIVE",    GREEN,    "42,800 records"),
    ("Salesforce CRM",       "Accounts · Opps · Cases · NPS",          "LIVE",    GREEN,    "18,400 records"),
    ("Workday HR",           "Headcount · Skills · Performance",        "LIVE",    GREEN,    "3,200 records"),
    ("Reuters Newsfeed",     "Market · Supply Chain · Regulatory",      "LIVE",    GREEN,    "1,240 signals"),
    ("Bloomberg Market",     "FX · Commodities · Indices",              "LIVE",    GREEN,    "580 tickers"),
    ("Microsoft Exchange",   "Email intelligence · Sentiment",          "LIVE",    GREEN,    "89,000 msgs"),
    ("Microsoft Teams",      "Meeting summaries · Actions",             "LIVE",    GREEN,    "480 meetings"),
    ("Plant IoT Sensors",    "OEE · Machine uptime · Cycle time",       "BETA",    AMBER,    "3 plants"),
    ("Competitor Monitor",   "Pricing · Launches · Patents · Jobs",     "LIVE",    GREEN,    "284 items"),
    ("Oracle Financials",    "GL · AR · AP · Budget vs Actual",         "PLANNED", GREY_TX,  "Q2 2026"),
    ("ServiceNow ITSM",      "Incidents · Change · Risk register",      "PLANNED", GREY_TX,  "Q2 2026"),
    ("Jira / Confluence",    "Projects · Delivery risk · Docs",         "PLANNED", GREY_TX,  "Q3 2026"),
]

cols = 3
row_h = 0.75
col_w = 4.2
start_y = 1.2
for idx, (name, detail, status, sc, records) in enumerate(connectors):
    col = idx % cols
    row = idx // cols
    cx = 0.3 + col * (col_w + 0.13)
    cy = start_y + row * (row_h + 0.08)
    bg = GREEN_LT if status == "LIVE" else (AMBER_LT if status == "BETA" else SLATE50)
    box(s, cx, cy, col_w, row_h, fill=WHITE, border=GREY_BD)
    # status dot box
    dot = s.shapes.add_shape(1, Inches(cx+0.1), Inches(cy+0.26), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = sc
    dot.line.fill.background()
    txt(s, name,    cx+0.3,  cy+0.06, col_w-0.5, 0.28, size=10, bold=True, color=SLATE1)
    txt(s, detail,  cx+0.3,  cy+0.33, col_w-0.8, 0.25, size=8,  color=GREY_TX)
    txt(s, records, cx+col_w-1.3, cy+0.06, 1.2, 0.28, size=9, bold=True, color=sc, align=PP_ALIGN.RIGHT)
    pill_label(s, status, cx+0.1, cy+0.49, bg, sc, size=7)

slide_number(s, 3)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ROLES
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "Role-Based Intelligence", "One platform. Three specialised experiences. Zero information overload.", AMBER)

roles = [
    {
        "title": "CEO / C-Suite",
        "sub":   "Strategic Command",
        "color": AMBER, "bg": AMBER_LT, "border": AMBER,
        "features": [
            "Insights Hub — KPI-driven morning view with AI narrative",
            "Analyst Intelligence — analyst insights surface directly",
            "Decision Hub — COMMIT / HOLD with evidence bundles",
            "Board Brief — 1-click auto-generated board pack",
            "Snap & Share — snapshot any view, download as PDF",
            "Execution Hub — track decisions through to outcome",
        ]
    },
    {
        "title": "Senior Analyst",
        "sub":   "Intelligence Bridge  ★ THE DIFFERENTIATOR",
        "color": VIOLET, "bg": VIOLET_LT, "border": VIOLET,
        "features": [
            "Analyst Studio — natural language over live data feeds",
            "AI Insight Generation — question → structured insight",
            "Share with CEO — push insight to CEO view. No email.",
            "Prompt Library — save & reuse intelligence templates",
            "Recall Control — recall shared insight any time",
            "Full Signals Access — annotate, verify, link to decisions",
        ]
    },
    {
        "title": "Platform Admin",
        "sub":   "System Command",
        "color": RED, "bg": RED_LT, "border": RED,
        "features": [
            "Data Connections — manage all 9+ MCP connectors",
            "System Prompts — configure & schedule AI auto-runs",
            "Audit Log — every action tracked with full history",
            "User Management — RBAC, roles, permissions",
            "Prompt Scheduler — 15min / hourly / daily / weekly",
            "Security Controls — data gating, ISO 27001 aligned",
        ]
    },
]

col_w = 4.0
for i, role in enumerate(roles):
    cx = 0.3 + i * (col_w + 0.27)
    box(s, cx, 1.18, col_w, 5.9, fill=role["bg"], border=role["border"], border_pt=1.5)
    # header band inside card
    hdr = s.shapes.add_shape(1, Inches(cx), Inches(1.18), Inches(col_w), Inches(0.72))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = role["color"]
    hdr.line.fill.background()
    txt(s, role["title"], cx+0.15, 1.2, col_w-0.3, 0.35, size=14, bold=True, color=WHITE)
    txt(s, role["sub"],   cx+0.15, 1.54, col_w-0.3, 0.28, size=8.5, color=WHITE)

    for j, feat in enumerate(role["features"]):
        fy = 2.12 + j * 0.75
        # bullet dot
        dot = s.shapes.add_shape(1, Inches(cx+0.18), Inches(fy+0.09), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = role["color"]
        dot.line.fill.background()
        txt(s, feat, cx+0.38, fy, col_w-0.5, 0.65, size=9.5, color=SLATE1)

slide_number(s, 4)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CEO MORNING VIEW (product mockup described)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "The CEO's Morning View", "Marcus Gaksh opens his laptop at 7am — here is what he sees", GREEN)

# Left panel — Analyst Intelligence + KPIs
box(s, 0.3, 1.2, 5.9, 5.85, fill=WHITE, border=GREY_BD)
txt(s, "ANALYST INTELLIGENCE", 0.45, 1.28, 5.6, 0.28, size=8, bold=True, color=AMBER)
box(s, 0.45, 1.58, 5.6, 1.35, fill=AMBER_LT, border=AMBER, border_pt=1)
txt(s, "Supply Chain Risk & Exposure Report", 0.6, 1.63, 5.3, 0.3, size=10, bold=True, color=AMBER)
txt(s, "7 active risks · $3.1M exposure · Titanium lead slip 18 days\nShared by Arjun Mehta · 09:14",
    0.6, 1.93, 5.3, 0.5, size=8.5, color=SLATE3)
pill_label(s, "NEW", 0.6, 2.42, AMBER_LT, AMBER, 8)
pill_label(s, "Pin to Board", 1.1, 2.42, BLUE_LT, BLUE, 8)

# KPI strip
kpis = [("$127.4M", "YTD Revenue", "+4.2%", GREEN),
        ("23.4%",   "EBITDA",       "-1.8pp", RED),
        ("84.2%",   "On-Time Del.", "-7.8pp", RED)]
for i, (val, lbl, delta, dc) in enumerate(kpis):
    kx = 0.45 + i * 1.9
    box(s, kx, 3.05, 1.78, 0.9, fill=BLUE_LT, border=BLUE, border_pt=0.75)
    txt(s, val,   kx+0.07, 3.1,  1.65, 0.3, size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, lbl,   kx+0.07, 3.42, 1.65, 0.22, size=7, color=GREY_TX, align=PP_ALIGN.CENTER)
    txt(s, delta, kx+0.07, 3.63, 1.65, 0.22, size=8, bold=True, color=dc, align=PP_ALIGN.CENTER)

# Decision card
box(s, 0.45, 4.1, 5.6, 1.82, fill=WHITE, border=GREY_BD)
txt(s, "DEC-001  ·  Supply Chain  ·  HIGH CONFIDENCE", 0.6, 4.15, 5.3, 0.25, size=7.5, color=GREY_TX)
txt(s, "Pre-buy 6-Month Titanium Supply at Current Rates", 0.6, 4.4, 5.3, 0.38, size=11, bold=True, color=SLATE1)
txt(s, "Bharat Forge lead time +18 days. Lock price before May increase.\nValue at stake: $1.1M",
    0.6, 4.79, 5.3, 0.45, size=8.5, color=SLATE3)
box(s, 0.6, 5.36, 1.1, 0.38, fill=GREEN, border=GREEN, border_pt=0)
txt(s, "✓ COMMIT", 0.6, 5.39, 1.1, 0.3, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 1.82, 5.36, 0.9, 0.38, fill=AMBER_LT, border=AMBER, border_pt=1)
txt(s, "⏸ HOLD",  1.82, 5.39, 0.9, 0.3, size=9, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Right panel — Insight + Signal + Risk
box(s, 6.45, 1.2, 6.55, 5.85, fill=WHITE, border=GREY_BD)

box(s, 6.6, 1.28, 6.25, 1.5, fill=GREEN_LT, border=GREEN, border_pt=1)
txt(s, "⚡ Revenue (KPI-001) — AI Analysis", 6.75, 1.32, 5.95, 0.28, size=9, bold=True, color=GREEN)
txt(s, "YTD Revenue $127.4M vs $135M target (94.4%). Gap of $7.6M driven by Bosch re-phasing ($4.1M) and Magna PO delay ($2.8M). Pipeline coverage 2.3× — healthy.",
    6.75, 1.6, 5.95, 0.68, size=8.5, color=SLATE3)
txt(s, "→ Prioritise CEO outreach to Bosch NA this week.", 6.75, 2.28, 5.95, 0.28, size=8.5, bold=True, color=BLUE)

box(s, 6.6, 2.9, 6.25, 1.38, fill=RED_LT, border=RED, border_pt=1)
txt(s, "🔴 HIGH  ·  Supply Chain Signal", 6.75, 2.95, 5.95, 0.25, size=8, bold=True, color=RED)
txt(s, "Chennai Port Congestion — 4–7 Day Transit Delays", 6.75, 3.2, 5.95, 0.3, size=10, bold=True, color=SLATE1)
txt(s, "3 GIS shipments affected. Expedite routing for PO-8821.", 6.75, 3.5, 5.95, 0.55, size=8.5, color=SLATE3)
txt(s, "Confidence 91%", 10.7, 2.95, 1.5, 0.25, size=8, bold=True, color=RED, align=PP_ALIGN.RIGHT)

box(s, 6.6, 4.43, 6.25, 1.48, fill=RGBColor(0xFF,0xF7,0xED), border=AMBER, border_pt=1)
txt(s, "⚠ DOWNSIDE RISK  ·  Regulatory", 6.75, 4.48, 5.95, 0.25, size=8, bold=True, color=AMBER)
txt(s, "EU CBAM Compliance — 32% Ready vs 80% Required Q4 2026", 6.75, 4.73, 5.95, 0.3, size=10, bold=True, color=SLATE1)
txt(s, "Non-compliance risks $24M Tier-1 contract renewals.", 6.75, 5.03, 5.95, 0.25, size=8.5, color=SLATE3)
box(s, 6.75, 5.43, 1.1, 0.32, fill=RED, border=RED, border_pt=0)
txt(s, "Escalate", 6.75, 5.45, 1.1, 0.26, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 7.98, 5.43, 1.3, 0.32, fill=BLUE_LT, border=BLUE, border_pt=0.75)
txt(s, "View Decision →", 7.98, 5.45, 1.3, 0.26, size=8, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

slide_number(s, 5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE DIFFERENTIATOR
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "The Market Differentiator — Analyst → CEO Intelligence Bridge",
            "No BI tool. No dashboard. No AI assistant. Nobody does this.", VIOLET)

steps = [
    ("1", "Analyst types\na question",
     '"What is our Q2 revenue risk?"',   VIOLET_LT, VIOLET),
    ("2", "AI synthesises\nlive data",
     "CRM + ERP + Meetings\n→ structured insight",            BLUE_LT, BLUE),
    ("3", "Analyst clicks\nShare with CEO",
     "One click. Instant.\nNo email. No PowerPoint.",          AMBER_LT, AMBER),
    ("4", "CEO sees it in\nmorning view",
     "Analyst Intelligence\nsection — with metrics",          RED_LT, RED),
    ("5", "CEO commits\nthe decision",
     "Full traceability\nin audit log",                        GREEN_LT, GREEN),
]
step_w = 2.35
for i, (num, title, detail, bg, fg) in enumerate(steps):
    sx = 0.3 + i*(step_w+0.15)
    box(s, sx, 1.25, step_w, 2.8, fill=bg, border=fg, border_pt=1.5)
    # number circle
    nc = s.shapes.add_shape(1, Inches(sx+step_w/2-0.22), Inches(1.3), Inches(0.44), Inches(0.44))
    nc.fill.solid(); nc.fill.fore_color.rgb = fg
    nc.line.fill.background()
    txt(s, num, sx+step_w/2-0.22, 1.3, 0.44, 0.44, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title,  sx+0.1, 1.82, step_w-0.2, 0.55, size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(s, detail, sx+0.1, 2.45, step_w-0.2, 0.85, size=9,  color=SLATE3, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        ax = sx + step_w + 0.01
        arr = s.shapes.add_shape(1, Inches(ax), Inches(2.5), Inches(0.12), Inches(0.12))
        arr.fill.solid(); arr.fill.fore_color.rgb = GREY_TX
        arr.line.fill.background()

# Bottom proof box
box(s, 0.3, 4.3, 12.73, 2.55, fill=WHITE, border=VIOLET, border_pt=1.5)
txt(s, "Why This Is Unique", 0.55, 4.38, 12, 0.35, size=13, bold=True, color=VIOLET)
proofs = [
    "Traditional flow: Analyst finds insight → writes report → emails CEO → CEO reads deck 3 days later → decides with stale data",
    "SiBoNi flow: Analyst asks question → AI generates insight from live data → one click → CEO sees it in next session → decides same day",
    "Result: 4.8× higher analyst insight utilisation · 73% faster executive decision cycle · Zero report-writing overhead",
]
for i, p in enumerate(proofs):
    col = [RED, GREEN, BLUE][i]
    dot2 = s.shapes.add_shape(1, Inches(0.5), Inches(4.82+i*0.62), Inches(0.12), Inches(0.12))
    dot2.fill.solid(); dot2.fill.fore_color.rgb = col; dot2.line.fill.background()
    txt(s, p, 0.72, 4.77+i*0.62, 12.1, 0.55, size=9.5, color=SLATE1 if i>0 else SLATE3)

slide_number(s, 6)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — COMPETITIVE TABLE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "Competitive Positioning", "vs Power BI · Tableau · ThoughtSpot · Glean AI · Palantir AIP", RED)

headers = ["Capability", "⚡ SiBoNi", "Power BI", "Tableau", "ThoughtSpot", "Glean AI", "Palantir AIP"]
col_xs  = [0.3, 4.3, 5.95, 7.1, 8.25, 9.45, 10.65]
col_ws  = [3.85, 1.5, 1.0, 1.0, 1.0, 1.0, 1.0]
rows = [
    ("Real-time multi-source data fusion",      ["✓","~","~","✓","✗","✓"]),
    ("AI-generated executive narrative",         ["✓","✗","✗","~","✓","~"]),
    ("Analyst → CEO direct push (no email) ★",   ["✓","✗","✗","✗","✗","✗"]),
    ("Natural language prompt over live data",   ["✓","✗","~","✓","✓","~"]),
    ("Decision tracking with evidence bundles",  ["✓","✗","✗","✗","✗","~"]),
    ("COMMIT / HOLD decision workflow",          ["✓","✗","✗","✗","✗","~"]),
    ("1-click board brief generation",           ["✓","✗","✗","✗","✗","✗"]),
    ("Admin-configurable prompt schedules",      ["✓","✗","✗","✗","~","~"]),
    ("KPI drill-down: signal / decision / risk", ["✓","~","~","~","✗","✓"]),
    ("Executive role-based access (no noise)",   ["✓","~","~","~","✗","~"]),
]

# Header row
row_h = 0.44
hy = 1.18
box(s, 0.3, hy, 12.73, row_h, fill=BLUE, border=BLUE, border_pt=0)
for j, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
    align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    txt(s, hdr, cx+0.06, hy+0.06, cw-0.1, row_h-0.1, size=9, bold=True, color=WHITE, align=align)

CHECK, CROSS, PART = "✓", "✗", "~"
c_col = {CHECK: GREEN, CROSS: RED, PART: AMBER}

for ri, (cap, vals) in enumerate(rows):
    ry = hy + (ri+1)*row_h
    bg = SLATE50 if ri % 2 == 0 else WHITE
    # highlight SiBoNi col
    box(s, col_xs[1], ry, col_ws[1], row_h, fill=BLUE_LT, border=BLUE, border_pt=0.5)
    box(s, col_xs[0], ry, col_ws[0], row_h, fill=bg, border=GREY_BD, border_pt=0.5)
    cap_color = VIOLET if "★" in cap else SLATE1
    cap_bold  = "★" in cap
    txt(s, cap, col_xs[0]+0.08, ry+0.07, col_ws[0]-0.14, row_h-0.1, size=8.5,
        bold=cap_bold, color=cap_color)
    for j, (v, cx, cw) in enumerate(zip(vals, col_xs[1:], col_ws[1:])):
        col = c_col.get(v, GREY_TX)
        if j == 0: col = GREEN if v == CHECK else col  # SiBoNi always green
        txt(s, v, cx+0.06, ry+0.07, cw-0.1, row_h-0.1, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

txt(s, "✓ Full  ·  ~ Partial / needs config  ·  ✗ Not available  ·  ★ Unique to SiBoNi",
    0.3, 7.08, 12.73, 0.28, size=8, color=GREY_TX)
slide_number(s, 7)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — IMPACT METRICS
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "Impact Metrics", "Measured outcomes from SiBoNi deployments", GREEN)

metrics = [
    ("73%",   "Reduction in time\nto executive decision",  BLUE),
    ("4.8×",  "Analyst insight\nutilisation by CEO",       VIOLET),
    ("1 day", "Time to first\nexecutive value",             GREEN),
    ("9+",    "Enterprise systems\nconnected at launch",   CYAN),
    ("$0",    "Additional BI team\nheadcount required",    AMBER),
]
mw = 2.35
for i, (val, lbl, col) in enumerate(metrics):
    mx = 0.3 + i*(mw+0.17)
    box(s, mx, 1.25, mw, 2.1, fill=WHITE, border=col, border_pt=1.5)
    txt(s, val, mx, 1.45, mw, 0.75, size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, lbl, mx+0.1, 2.25, mw-0.2, 0.7, size=10, color=GREY_TX, align=PP_ALIGN.CENTER)

# Quote
box(s, 0.5, 3.65, 12.33, 1.42, fill=BLUE_LT, border=BLUE, border_pt=1.5)
txt(s,
    '"The gap between data and decisions is not a technology problem. It\'s a communication problem. SiBoNi solves it."',
    0.8, 3.8, 11.8, 0.75, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER, italic=True)
txt(s, "— SiBoNi Product Philosophy", 0.8, 4.6, 11.8, 0.3, size=9, color=GREY_TX, align=PP_ALIGN.CENTER)

# Supporting stats
stats = [
    ("3 days → 60 min", "Revenue anomaly detection"),
    ("2 weeks → 2 hrs",  "Competitor response cycle"),
    ("2 days → 4 min",   "Board pack generation"),
    ("3 weeks early",    "Workforce capacity decisions"),
]
sw = 3.0
for i, (stat, lbl) in enumerate(stats):
    sx = 0.3 + i*(sw+0.1)
    box(s, sx, 5.32, sw, 0.88, fill=GREEN_LT, border=GREEN, border_pt=0.75)
    txt(s, stat, sx+0.1, 5.37, sw-0.2, 0.35, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, lbl,  sx+0.1, 5.73, sw-0.2, 0.35, size=8.5, color=SLATE3, align=PP_ALIGN.CENTER)

slide_number(s, 8)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — USE CASES
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "Real-World Use Cases", "Six scenarios that demonstrate measurable executive value", CYAN)

use_cases = [
    ("1", "Morning Revenue Gap Alert",
     "7:00am — System Prompt 'Revenue Anomaly Scan' runs against ERP + CRM. Detects $2.4M shortfall. CEO sees it before first meeting. Calls Magna by 8am.",
     "⏱ 3 days → 60 minutes", GREEN),
    ("2", "Supply Chain Pre-Buy Authorisation",
     "Analyst detects Titanium lead-time slippage. Shares insight to CEO. CEO clicks COMMIT on DEC-001 with evidence. Procurement authorised same day.",
     "⚡ $980K decision in 4 hours — averted $1.1M loss", AMBER),
    ("3", "Board Meeting Preparation",
     "CFO asks for board pack. CEO opens Board Brief — auto-generated from live KPIs, decisions, risks, pinned analyst insights. Clicks Export PDF.",
     "📋 2 days → 4 minutes", BLUE),
    ("4", "Competitor Pricing Response",
     "Competitor Monitor detects Bosch 8% price cut. System Prompt fires. CEO sees signal. Requests analyst to model rebate strategy. Shared back same hour.",
     "🔍 2-week research cycle → 2 hours", VIOLET),
    ("5", "Workforce Capacity Planning",
     "HR + ERP cross-reference: Detroit at 78% vs 85% threshold. Analyst shares insight: 14 machinists needed by Apr 20, $2.1M overtime risk. CEO authorises staffing.",
     "👥 Staffing decision 3 weeks earlier", RED),
    ("6", "Regulatory Compliance Escalation",
     "Risk monitoring: EU CBAM readiness 32% vs 80% required. Analyst escalates. CEO activates $240K compliance programme. Full audit trail in one system.",
     "🏛 $24M contract risk addressed 9 months early", CYAN),
]

uc_w = 6.1
uc_h = 1.72
for idx, (num, title, body, result, col) in enumerate(use_cases):
    col_pos = idx % 2
    row_pos = idx // 2
    ux = 0.3 + col_pos*(uc_w+0.2)
    uy = 1.22 + row_pos*(uc_h+0.12)
    box(s, ux, uy, uc_w, uc_h, fill=WHITE, border=GREY_BD)
    # num badge
    nb = s.shapes.add_shape(1, Inches(ux+0.12), Inches(uy+0.12), Inches(0.38), Inches(0.38))
    nb.fill.solid(); nb.fill.fore_color.rgb = col; nb.line.fill.background()
    txt(s, num, ux+0.12, uy+0.12, 0.38, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, ux+0.65, uy+0.1,  uc_w-0.77, 0.3, size=10, bold=True, color=SLATE1)
    txt(s, body,  ux+0.65, uy+0.42, uc_w-0.77, 0.72, size=8.5, color=SLATE3)
    # result bar
    box(s, ux+0.1, uy+uc_h-0.44, uc_w-0.2, 0.33, fill=GREEN_LT, border=GREEN, border_pt=0.5)
    txt(s, result, ux+0.22, uy+uc_h-0.41, uc_w-0.4, 0.28, size=8.5, bold=True, color=GREEN)

slide_number(s, 9)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, OFF_W)
header_band(s, "Technical Architecture", "Built for enterprise. Fast to deploy. Cloud-agnostic.", SLATE3)

layers = [
    ("Presentation Layer", "React 18 + Vite · Role-based SPA",
     "CEO Cockpit  ·  Analyst Studio  ·  Admin Panel  ·  Insights Hub  ·  Decision Hub  ·  Execution Hub  ·  Board Brief  ·  Signals Console",
     BLUE, BLUE_LT),
    ("AI Intelligence Layer", "Decision Engine · Prompt Orchestrator · Confidence Scoring",
     "System Prompts (scheduled)  ·  Analyst Prompt Builder  ·  Cross-source Fusion Engine  ·  Anomaly Detection  ·  KPI Computation  ·  Signal Scoring",
     VIOLET, VIOLET_LT),
    ("State & Context Layer", "React Context · localStorage · Real-time sync",
     "FilterContext  ·  DecisionStateContext  ·  PromptContext  ·  AuthContext  ·  WorkspaceContext  ·  Full Audit Trail",
     AMBER, AMBER_LT),
    ("MCP Connector Layer", "Model Context Protocol · Unified Schema · Auth Governance",
     "SAP ERP  ·  Salesforce CRM  ·  Workday HR  ·  Reuters  ·  Bloomberg  ·  Exchange  ·  Teams  ·  IoT Sensors  ·  Competitor Feed",
     GREEN, GREEN_LT),
    ("Infrastructure", "Cloud-agnostic · Azure / AWS / GCP · SOC 2 Type II · ISO 27001",
     "RBAC Auth  ·  Encryption at rest & transit  ·  Full Audit Log  ·  API Gateway  ·  Edge CDN  ·  Docker / Kubernetes",
     SLATE3, SLATE50),
]

ly = 1.2
for i, (lname, lsub, lchips, col, bg) in enumerate(layers):
    lh = 0.92
    box(s, 0.3, ly, 12.73, lh, fill=bg, border=col, border_pt=1.5)
    # left colour tab
    tab = s.shapes.add_shape(1, Inches(0.3), Inches(ly), Inches(0.12), Inches(lh))
    tab.fill.solid(); tab.fill.fore_color.rgb = col; tab.line.fill.background()
    txt(s, lname, 0.56, ly+0.05, 3.8, 0.3, size=10, bold=True, color=col)
    txt(s, lsub,  0.56, ly+0.37, 4.0, 0.28, size=8, color=GREY_TX)
    txt(s, lchips, 4.6, ly+0.18, 8.3, 0.6, size=8.5, color=SLATE1)
    ly += lh + 0.08

# Bottom SLA strip
slas = [("<2s AI latency", BLUE), ("99.9% SLA", GREEN), ("SOC 2 Type II", VIOLET), ("1-day deploy", AMBER), ("ISO 27001", RED)]
for i, (lbl, col) in enumerate(slas):
    sx = 0.3 + i*2.55
    box(s, sx, 6.55, 2.35, 0.62, fill=WHITE, border=col, border_pt=1)
    txt(s, lbl, sx+0.1, 6.65, 2.15, 0.38, size=10, bold=True, color=col, align=PP_ALIGN.CENTER)

slide_number(s, 10)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CTA
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s, WHITE)

bar2 = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
bar2.fill.solid(); bar2.fill.fore_color.rgb = BLUE; bar2.line.fill.background()

txt(s, "SiBoNi  CXO Cockpit", 0.5, 0.2, 12.33, 0.5, size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txt(s, "Your CEO's morning view starts here.", 0.5, 0.82, 12.33, 0.45, size=18, bold=True, color=SLATE1, align=PP_ALIGN.CENTER)
txt(s, "Connect your ERP and CRM in under a day. Actionable intelligence by tomorrow morning.",
    0.5, 1.35, 12.33, 0.38, size=12, color=GREY_TX, align=PP_ALIGN.CENTER)

# 3 option boxes
options = [
    ("🚀 Pilot in 1 Day",    "Connect 2 data sources.\nConfigure first system prompt.\nCEO dashboard live by EOD.", "Free 30 days", GREEN),
    ("🎯 Live Demo",         "Watch Analyst Studio live:\ntype → share → CEO sees it.\n45-minute session.", "Book a session", BLUE),
    ("📋 Enterprise Deal",   "Full deployment across all roles.\nCustom MCP connectors.\nDedicated success team.", "Talk to us", VIOLET),
]
ow = 3.8
for i, (title, body, cta_txt, col) in enumerate(options):
    ox = 0.5 + i*(ow+0.35)
    box(s, ox, 1.95, ow, 2.85, fill=WHITE, border=col, border_pt=1.5)
    hb = s.shapes.add_shape(1, Inches(ox), Inches(1.95), Inches(ow), Inches(0.52))
    hb.fill.solid(); hb.fill.fore_color.rgb = col; hb.line.fill.background()
    txt(s, title, ox+0.12, 1.98, ow-0.24, 0.42, size=12, bold=True, color=WHITE)
    txt(s, body,  ox+0.12, 2.6,  ow-0.24, 1.35, size=10, color=SLATE1)
    box(s, ox+0.15, 4.35, ow-0.3, 0.35, fill=col, border=col, border_pt=0)
    txt(s, cta_txt, ox+0.15, 4.37, ow-0.3, 0.3, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Demo credentials
box(s, 0.5, 5.02, 12.33, 1.78, fill=SLATE50, border=GREY_BD)
txt(s, "Try it now — demo accounts", 0.7, 5.1, 11, 0.3, size=10, bold=True, color=SLATE1)
creds = [
    ("👑 CEO View",       "ceo@gis.com",      "ceo123",      AMBER),
    ("🔬 Analyst Studio", "analyst@gis.com",  "analyst123",  VIOLET),
    ("⚙️ Admin Panel",   "admin@gis.com",    "admin123",    RED),
]
for i, (role, email, pw, col) in enumerate(creds):
    cx = 0.6 + i*4.15
    box(s, cx, 5.42, 3.85, 1.22, fill=WHITE, border=col, border_pt=1)
    txt(s, role,  cx+0.12, 5.46, 3.6, 0.3, size=10, bold=True, color=col)
    txt(s, email, cx+0.12, 5.77, 3.6, 0.28, size=9, color=GREY_TX)
    txt(s, pw,    cx+0.12, 6.03, 3.6, 0.28, size=9, color=SLATE3)

# Trust badges
badges = ["🔒 SOC 2 Type II","🔐 ISO 27001","☁️ Cloud-agnostic","⚡ 1-day deploy","🌍 Multi-region","🏗️ MCP Standard"]
txt(s, "  ·  ".join(badges), 0.5, 7.08, 12.33, 0.3, size=8, color=GREY_TX, align=PP_ALIGN.CENTER)

bar3 = s.shapes.add_shape(1, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12))
bar3.fill.solid(); bar3.fill.fore_color.rgb = BLUE; bar3.line.fill.background()

slide_number(s, 11)

# ════════════════════════════════════════════════════════════════════════════
out = "siboni-deck.pptx"
prs.save(out)
print(f"Saved: {out} ({prs.slides.__len__()} slides)")
